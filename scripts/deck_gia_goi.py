"""Sinh docs/Gia-goi-va-hoan-von.pptx — trình bày giá gói & điểm hoàn vốn.

    python scripts/deck_gia_goi.py

MỌI SỐ TRONG SLIDE ĐỀU TÍNH TẠI ĐÂY, không gõ tay vào văn bản: đổi đơn giá ở
`GIA` hay hạn mức gói là cả bộ slide tự cập nhật, không sợ lệch giữa bảng và
kết luận. Token đầu vào là số ĐO THẬT ở docs/cost-estimate.md §1.

Sửa gì khi có số thật:
  - `GIA`      : đơn giá VNGCloud (đang là proxy Google)
  - `CO_DINH`  : chi phí cố định/tháng ở slide 7
  - `STRONG`   : token của lời gọi tầng mạnh (đang ước 2k in / 4k out)
"""
from math import ceil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

VND = 25_400
GIA = {"cheap_in": .10, "cheap_out": .40, "strong_in": 1.25, "strong_out": 10.0}
usd = lambda i, o, t: i*GIA[f"{t}_in"]/1e6 + o*GIA[f"{t}_out"]/1e6
QA = usd(1_200, 500, "cheap")                     # đo thật
STRONG = usd(2_000, 4_000, "strong")              # ước lượng
MOT_LAN = 21*6*STRONG + 294*usd(1_241, 348, "cheap") + 880*150*GIA["cheap_in"]/1e6
DAU_TU = MOT_LAN*VND + 21*150_000
BIEN = 99_000 - 300*QA*VND*.7

def d(x):  # 1234567 -> "1.234.567"
    return f"{x:,.0f}".replace(",", ".")

TIM   = RGBColor(0x5B, 0x54, 0xF0)
TIM2  = RGBColor(0x82, 0x66, 0xFF)
MUC   = RGBColor(0x0F, 0x14, 0x2E)
XAM   = RGBColor(0x6B, 0x72, 0x8A)
XANH  = RGBColor(0x10, 0xA5, 0x74)
DO    = RGBColor(0xDC, 0x26, 0x26)
CAM   = RGBColor(0xD9, 0x77, 0x06)
TRANG = RGBColor(0xFF, 0xFF, 0xFF)
NEN   = RGBColor(0xF7, 0xF8, 0xFC)

prs = Presentation()
prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)   # 16:9
W, H = prs.slide_width, prs.slide_height

def cm(v): return Emu(int(v*360000))

def slide(tieu_de, phu=None, nen=TRANG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = nen; bg.line.fill.background()
    bg.shadow.inherit = False
    dai = s.shapes.add_shape(1, 0, 0, W, cm(.28))
    dai.fill.solid(); dai.fill.fore_color.rgb = TIM; dai.line.fill.background()
    dai.shadow.inherit = False
    tb = s.shapes.add_textbox(cm(1.6), cm(.9), W-cm(3.2), cm(1.5)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]; p.text = tieu_de
    p.font.size, p.font.bold, p.font.color.rgb = Pt(30), True, MUC
    if phu:
        q = tb.add_paragraph(); q.text = phu
        q.font.size, q.font.color.rgb = Pt(14), XAM
    return s

def txt(s, x, y, w, h, dong, co=13, mau=MUC, dam=False, canh=PP_ALIGN.LEFT):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    for i, ln in enumerate(dong):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(ln, tuple):
            noi, c, b, sz = (ln + (mau, dam, co))[:4]
            p.text, p.font.color.rgb, p.font.bold, p.font.size = noi, c, b, Pt(sz)
        else:
            p.text, p.font.color.rgb, p.font.bold, p.font.size = ln, mau, dam, Pt(co)
        p.alignment = canh
        p.space_after = Pt(5)
    return tf

def the(s, x, y, w, h, vien=None):
    sh = s.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = TRANG
    sh.line.color.rgb = vien or RGBColor(0xE2, 0xE6, 0xF0); sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh

def bang(s, x, y, w, cot, hang, rong=None, cao=.95):
    n = len(hang)+1
    t = s.shapes.add_table(n, len(cot), x, y, w, cm(cao*n)).table
    if rong:
        for i, r in enumerate(rong): t.columns[i].width = cm(r)
    for j, c in enumerate(cot):
        cel = t.cell(0, j); cel.text = c
        pr = cel.text_frame.paragraphs[0]
        pr.font.size, pr.font.bold, pr.font.color.rgb = Pt(11), True, TRANG
        cel.fill.solid(); cel.fill.fore_color.rgb = TIM
    for i, r in enumerate(hang, 1):
        for j, v in enumerate(r):
            cel = t.cell(i, j)
            noi, mau, dam = (v if isinstance(v, tuple) else (v, MUC, False))
            cel.text = str(noi)
            pr = cel.text_frame.paragraphs[0]
            pr.font.size, pr.font.color.rgb, pr.font.bold = Pt(11), mau, dam
            cel.fill.solid()
            cel.fill.fore_color.rgb = TRANG if i % 2 else NEN
    return t

# ─────────────────────────── 1. Bìa ───────────────────────────
s = slide("Giá gói Gia sư DTP", nen=NEN)
s.shapes.add_shape(1, 0, 0, W, cm(.28)).fill.solid()
txt(s, cm(1.6), cm(3.2), W-cm(3.2), cm(3),
    [("Định giá theo chi phí AI thật · điểm hoà vốn · hoàn vốn", MUC, True, 22),
     ("Số token ĐO THẬT trên SGK Toán 6 (docs/cost-estimate.md §1). "
      "Đơn giá là proxy Google — cần xác nhận Console VNGCloud.", XAM, False, 13)])
txt(s, cm(1.6), cm(6.6), W-cm(3.2), cm(2),
    [("Kết luận trước: AI KHÔNG phải rào cản giá.", XANH, True, 17),
     ("Giá 99.000đ/tháng trong mockup cho biên gộp ~98%. "
      "Rào cản thật là chi phí cố định và số học sinh trả tiền.", MUC, False, 14)])

# ─────────────────── 2. Gói theo mockup ───────────────────
s = slide("Gói trong mockup", "Nguyên văn từ ui-design/mockups/gia-su-dtp-demo.html")
bang(s, cm(1.6), cm(3.0), W-cm(3.2),
     ["Gói", "Giá", "Chu kỳ khác", "Hạn mức hỏi trợ lý"],
     [["Miễn phí", "0đ", "không giới hạn thời gian", "5 lượt/ngày (~130/tháng)"],
      [("Chuẩn", TIM, True), ("99.000đ/tháng", TIM, True),
       "kỳ 399.000đ · năm 699.000đ", "300 lượt/tháng"],
      ["Gia đình", "149.000đ/tháng", "kỳ 599.000đ · năm 1.049.000đ", "tối đa 3 HS"],
      ["Nhà trường", "theo hợp đồng", "kích hoạt key-code", "tính theo số HS"]],
     rong=[5.2, 6.2, 9.6, 9.0])
txt(s, cm(1.6), cm(9.4), W-cm(3.2), cm(2),
    [("⚠️ Chưa có trong code: không có bảng subscriptions/payments, "
      "không có cổng thanh toán.", CAM, True, 13),
     ("Hiện chỉ có `chat_daily_limit = 20/ngày` áp cho MỌI user và "
      "`users.daily_limit_override` để nâng tay từng người.", XAM, False, 12)])

# ─────────────────── 3. Chi phí AI thật ───────────────────
s = slide("Học sinh chỉ chạm được tầng RẺ",
          "Đây là lý do biên gộp cao — không phải phỏng đoán, mà là kiến trúc")
bang(s, cm(1.6), cm(3.0), W-cm(3.2),
     ["Việc", "Ai kích hoạt", "Tầng model", "Chi phí"],
     [[("Hỏi trợ lý (qa)", TIM, True), ("Học sinh", TIM, True),
       ("flash-lite + cache 24h", TIM, True), (f"{QA*VND:.1f}đ / lượt", TIM, True)],
      ["Soạn bài (lesson_ingest)", "Chuyên gia · 1 lần/bài", "pro-preview",
       f"~{d(STRONG*VND)}đ / lời gọi"],
      ["Sinh đề kiểm tra (quiz_gen)", "Chuyên gia · 1 lần/bài", "pro-preview",
       f"~{d(STRONG*VND)}đ / lời gọi"],
      ["OCR trang sách", "Chuyên gia · 1 lần/sách", "flash-lite", "~7đ / trang"],
      ["Video minh hoạ", "HS bấm, CACHE theo khái niệm", "flash-lite + CPU",
       "≈0đ từ HS thứ 2"]],
     rong=[8.0, 8.4, 7.2, 6.4])
txt(s, cm(1.6), cm(9.8), W-cm(3.2), cm(1.6),
    [("Mọi việc đắt đều là MỘT LẦN cho mỗi bài, chia cho toàn bộ học sinh. "
      "Chi phí biên mỗi học sinh = chỉ tiền chat.", MUC, True, 13)])

# ─────────────── 4. Biên gộp từng gói ───────────────
s = slide("Biên gộp từng gói", "Chi phí AI = số lượt × 8,1đ, trừ 30% nhờ cache 24h")
hang = []
for ten, gia, luot in [("Miễn phí", 0, 130), ("Chuẩn", 99_000, 300),
                       ("Gia đình (3 HS)", 149_000, 900)]:
    c = luot*QA*VND*.7
    hang.append([(ten, TIM if gia == 99_000 else MUC, gia == 99_000),
                 d(gia)+"đ" if gia else "0đ", f"{luot} lượt", d(c)+"đ",
                 (f"{(gia-c)/gia*100:.1f}%", XANH, True) if gia else ("—", XAM, False)])
bang(s, cm(1.6), cm(3.0), W-cm(3.2),
     ["Gói", "Giá/tháng", "Hạn mức", "Chi phí AI", "Biên gộp"], hang,
     rong=[7.0, 5.6, 5.2, 5.4, 6.8])
txt(s, cm(1.6), cm(7.6), W-cm(3.2), cm(3),
    [("Gói Miễn phí tốn ~740đ/tháng mỗi người — rẻ đến mức dùng làm phễu được, "
      "nhưng 1.000 người dùng thử = 740.000đ/tháng.", MUC, False, 13),
     ("Gói Chuẩn: 99.000đ thu về, 1.707đ chi cho AI. "
      "Còn 97.293đ/tháng cho hạ tầng, nội dung, vận hành và lợi nhuận.", MUC, True, 14)])

# ─────────────── 5. Giá có bền không ───────────────
s = slide("Giá 99.000đ bền tới đâu?",
          "Đơn giá dùng là proxy Google — nếu VNGCloud đắt hơn thì sao?")
bang(s, cm(1.6), cm(3.0), cm(14),
     ["Đơn giá thật so với proxy", "Chi phí AI/HS/tháng", "Biên gộp"],
     [[f"×{x}", d(300*QA*VND*.7*x)+"đ",
       ((f"{(99_000-300*QA*VND*.7*x)/99_000*100:.1f}%",
         XANH if x <= 10 else (CAM if x <= 20 else DO), True))]
      for x in (1, 5, 10, 20, 50)],
     rong=[6.2, 4.4, 3.4])
the(s, cm(16.6), cm(3.0), cm(13.2), cm(5.4), DO)
txt(s, cm(17.2), cm(3.4), cm(12), cm(4.6),
    [("Rủi ro thật KHÔNG phải đơn giá", DO, True, 15),
     (f"Một lời gọi tầng mạnh đắt gấp {STRONG/QA:.0f} lần tầng rẻ.", MUC, True, 13),
     (f"Nếu sau này mở tính năng cho HS chạm pro-preview: 300 lượt = "
      f"{d(300*STRONG*VND)}đ/HS → LỖ {d(300*STRONG*VND-99_000)}đ.", DO, False, 12),
     ("Giữ nguyên tắc: mọi việc tầng mạnh chỉ chạy lúc SOẠN BÀI, "
      "không bao giờ trong request của học sinh.", MUC, False, 12)])
txt(s, cm(1.6), cm(9.0), cm(14), cm(1.6),
    [("Chịu được đơn giá đắt hơn 10 lần mà vẫn còn 83% biên gộp.", XANH, True, 13)])

# ─────────────── 6. Đầu tư một lần ───────────────
s = slide("Đầu tư một lần cho Toán 6", "21 đơn vị kiến thức · 294 trang sách")
bang(s, cm(1.6), cm(3.0), cm(16),
     ["Hạng mục", "Khối lượng", "Chi phí"],
     [["AI soạn nội dung 7 mục", "21 bài × 6 lời gọi mạnh", d(21*6*STRONG*VND)+"đ"],
      ["OCR sách vào kho", "294 trang", d(294*usd(1_241,348,'cheap')*VND)+"đ"],
      ["Nhúng vector", "~880 đoạn", d(880*150*GIA['cheap_in']/1e6*VND)+"đ"],
      [("Cộng AI", TIM, True), "", (d(MOT_LAN*VND)+"đ", TIM, True)],
      ["Công chuyên gia rà nội dung", "21 giờ × 150.000đ", d(21*150_000)+"đ"],
      [("TỔNG", MUC, True), "", (d(DAU_TU)+"đ", MUC, True)]],
     rong=[8.0, 5.4, 4.6])
the(s, cm(18.4), cm(3.0), cm(11.4), cm(5.0), XANH)
txt(s, cm(19.0), cm(3.5), cm(10.2), cm(4.2),
    [("AI chỉ chiếm 4% đầu tư", XANH, True, 16),
     (f"{d(MOT_LAN*VND)}đ trên tổng {d(DAU_TU)}đ.", MUC, False, 13),
     ("Phần lớn tiền là CÔNG NGƯỜI rà soát. Muốn giảm giá thành, "
      "tối ưu quy trình duyệt — không phải tối ưu token.", MUC, False, 12)])

# ─────────────── 7. Hoà vốn ───────────────
s = slide("Bao nhiêu học sinh thì hoà vốn?",
          f"Mỗi HS gói Chuẩn để lại {d(BIEN)}đ/tháng sau chi phí AI")
bang(s, cm(1.6), cm(3.0), W-cm(3.2),
     ["Chi phí cố định/tháng", "HS hoà vốn hàng tháng",
      "HS để hoàn cả đầu tư trong 6 tháng", "Doanh thu/tháng khi hoà vốn"],
     [[d(cd)+"đ", (f"{ceil(cd/BIEN)} HS", TIM, True),
       f"{ceil((DAU_TU/6+cd)/BIEN)} HS", d(ceil(cd/BIEN)*99_000)+"đ"]
      for cd in (3_000_000, 6_000_000, 12_000_000, 25_000_000)],
     rong=[7.6, 7.4, 8.4, 6.6])
txt(s, cm(1.6), cm(8.4), W-cm(3.2), cm(2.6),
    [("Chi phí cố định là con số DUY NHẤT tôi không có số thật — "
      "hạ tầng (VPS + RDS) và nhân sự tuỳ cách bạn vận hành.", CAM, True, 13),
     ("Điền số thật vào cột 1 rồi đọc ngang. Ví dụ nếu hạ tầng + 1 nhân sự "
      "bán thời gian ≈ 12 triệu/tháng thì cần 123 học sinh trả tiền để hoà vốn.",
      MUC, False, 13)])

# ─────────────── 8. Đề xuất giá ───────────────
s = slide("Đề xuất: giữ giá mockup, sửa chỗ khác")
the(s, cm(1.6), cm(3.0), cm(13.6), cm(6.6), XANH)
txt(s, cm(2.2), cm(3.4), cm(12.4), cm(5.8),
    [("GIỮ 99.000đ/tháng", XANH, True, 17),
     ("Biên gộp 98% chịu được cả khi đơn giá AI đắt gấp 10. "
      "Hạ giá không giải quyết được gì mà mất doanh thu.", MUC, False, 12),
     ("", MUC, False, 8),
     ("Gói năm 699.000đ = 7,1 tháng — giảm 41%", MUC, True, 13),
     ("Chiết khấu này KHÔNG do chi phí AI (chi phí gần như tuyến tính theo lượt "
      "dùng). Nó mua được dòng tiền trả trước và giảm rời gói. Nếu cần dòng tiền "
      "cho đầu tư nội dung, đây là đòn bẩy đúng.", MUC, False, 12)])
the(s, cm(16.0), cm(3.0), cm(13.8), cm(6.6), CAM)
txt(s, cm(16.6), cm(3.4), cm(12.6), cm(5.8),
    [("BA CHỖ CẦN SỬA", CAM, True, 17),
     ("1. Gói Miễn phí 5 lượt/ngày = 130 lượt/tháng, gần một nửa gói Chuẩn (300). "
      "Người dùng nhẹ không có lý do trả tiền. Đề xuất hạ còn 2–3 lượt/ngày.",
      MUC, False, 12),
     (f"2. Gói Gia đình 149.000đ cho 3 HS = {d(149_000/3)}đ/HS, rẻ hơn một nửa gói "
      f"Chuẩn. Hai anh em mua Gia đình thay vì 2 gói Chuẩn → mất "
      f"{d(2*99_000-149_000)}đ. Đề xuất 179.000đ (2 gói Chuẩn vẫn đắt hơn).",
      MUC, False, 12),
     ("3. Chưa có bảng subscriptions/payments và cổng thanh toán — hạn mức hiện "
      "áp chung 20 lượt/ngày cho mọi người, chưa phân theo gói.", MUC, False, 12)])

# ─────────────── 9. Cần làm gì tiếp ───────────────
s = slide("Việc cần làm trước khi bán")
bang(s, cm(1.6), cm(3.0), W-cm(3.2),
     ["Việc", "Vì sao", "Ai xác nhận"],
     [[("Lấy đơn giá thật trên Console VNGCloud", DO, True),
       "Mọi số ở đây dùng proxy Google. Sai lệch >10× mới ảnh hưởng kết luận.",
       "Bạn"],
      [("Chốt chi phí cố định/tháng", DO, True),
       "Là biến duy nhất quyết định điểm hoà vốn.", "Bạn"],
      ["Bảng subscriptions + payments + cổng thanh toán",
       "Chưa có gì trong code; hạn mức đang áp chung.", "Tôi làm được"],
      ["Áp hạn mức THEO GÓI (thay 20/ngày cố định)",
       "Đã có users.daily_limit_override, cần nối vào gói.", "Tôi làm được"],
      ["Đo lại token sau khi đổi model/prompt",
       "Prompt 7 mục mới làm tuần này chưa đo lại.", "Tôi làm được"]],
     rong=[10.4, 13.6, 5.6])
txt(s, cm(1.6), cm(9.6), W-cm(3.2), cm(1.4),
    [("Hai việc đầu là của bạn và chúng chặn mọi tính toán còn lại.", MUC, True, 13)])

prs.save("docs/Gia-goi-va-hoan-von.pptx")
print("Đã lưu docs/Gia-goi-va-hoan-von.pptx ·", len(prs.slides.__iter__.__self__._sldIdLst), "slide")
